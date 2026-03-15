"""
gen_ppt_dinamico.py — Motor de generación de PPT dinámicas para SAC
===================================================================
Genera presentaciones PowerPoint con python-pptx (Python puro).
Incluye: métricas, pipeline SAC, gráficos, tablas, separadores.
Filtros: geográfico, tipo siniestro, empresa, rango de fechas.
"""

import io
import os
import json
import tempfile
import pandas as pd
import numpy as np
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ══════════════════════════════════════════════════════════════════
# FUNCIONES DE FILTRADO Y CÁLCULO (usadas por app.py para preview)
# ══════════════════════════════════════════════════════════════════

def _safe_col(df, col):
    return col in df.columns


def _aplicar_filtros(df, filtros):
    """Aplica todos los filtros al DataFrame."""
    df = df.copy()
    empresa = filtros.get("empresa", "ambas")
    if empresa != "ambas" and _safe_col(df, "EMPRESA"):
        df = df[df["EMPRESA"].astype(str).str.upper().str.contains(empresa.upper())]
    tipos = filtros.get("tipos_siniestro", [])
    if tipos and _safe_col(df, "TIPO_SINIESTRO"):
        df = df[df["TIPO_SINIESTRO"].isin(tipos)]
    fecha_inicio = filtros.get("fecha_inicio")
    fecha_fin = filtros.get("fecha_fin")
    col_fecha = filtros.get("col_fecha", "FECHA_AVISO")
    if fecha_inicio and fecha_fin and _safe_col(df, col_fecha):
        df[col_fecha] = pd.to_datetime(df[col_fecha], errors="coerce")
        mask = (df[col_fecha] >= pd.Timestamp(fecha_inicio)) & (df[col_fecha] <= pd.Timestamp(fecha_fin))
        df = df[mask]
    deptos = filtros.get("departamentos", [])
    provs = filtros.get("provincias", [])
    dists = filtros.get("distritos", [])
    if deptos and _safe_col(df, "DEPARTAMENTO"):
        df = df[df["DEPARTAMENTO"].isin(deptos)]
    if provs and _safe_col(df, "PROVINCIA"):
        df = df[df["PROVINCIA"].isin(provs)]
    if dists and _safe_col(df, "DISTRITO"):
        df = df[df["DISTRITO"].isin(dists)]
    return df


def _calcular_metricas(df):
    """Calcula métricas principales del DataFrame."""
    n = len(df)
    cerrados = 0
    if _safe_col(df, "ESTADO_INSPECCION"):
        cerrados = int(len(df[df["ESTADO_INSPECCION"].astype(str).str.upper() == "CERRADO"]))
    pct_eval = (cerrados / n * 100) if n > 0 else 0
    indem = float(df["INDEMNIZACION"].sum()) if _safe_col(df, "INDEMNIZACION") else 0
    desemb = float(df["MONTO_DESEMBOLSADO"].sum()) if _safe_col(df, "MONTO_DESEMBOLSADO") else 0
    pct_desemb = (desemb / indem * 100) if indem > 0 else 0
    ha = float(df["SUP_INDEMNIZADA"].sum()) if _safe_col(df, "SUP_INDEMNIZADA") else 0
    if _safe_col(df, "N_PRODUCTORES") and _safe_col(df, "INDEMNIZACION"):
        _indemn = pd.to_numeric(df["INDEMNIZACION"], errors="coerce").fillna(0)
        _prods = pd.to_numeric(df["N_PRODUCTORES"], errors="coerce").fillna(0)
        productores = int(_prods[_indemn > 0].sum())
    elif _safe_col(df, "N_PRODUCTORES"):
        productores = int(df["N_PRODUCTORES"].sum())
    else:
        productores = 0
    return {
        "avisos": n, "cerrados": cerrados, "pct_eval": round(pct_eval, 1),
        "indemnizacion": indem, "desembolso": desemb,
        "pct_desembolso": round(pct_desemb, 1),
        "ha_indemnizadas": round(ha, 2), "productores": productores,
    }


def _calcular_pipeline(df):
    """Calcula pipeline del proceso SAC."""
    if not _safe_col(df, "ESTADO_INSPECCION"):
        return []
    raw = df["ESTADO_INSPECCION"].astype(str).str.upper().value_counts().to_dict()
    order = ["NOTIFICADO", "PROGRAMADO", "REPROGRAMADO", "CERRADO"]
    result = []
    for estado in order:
        val = raw.get(estado, 0)
        if val > 0:
            result.append({"label": estado.title(), "val": int(val)})
    for k, v in raw.items():
        if k not in order and v > 0:
            result.append({"label": k.title(), "val": int(v)})
    return result


def _top_breakdown(df, col, n=10):
    """Top N por columna geográfica."""
    if not _safe_col(df, col):
        return []
    agg = {"Avisos": (col, "count")}
    if _safe_col(df, "INDEMNIZACION"):
        agg["Indemnización"] = ("INDEMNIZACION", "sum")
    if _safe_col(df, "MONTO_DESEMBOLSADO"):
        agg["Desembolso"] = ("MONTO_DESEMBOLSADO", "sum")
    if _safe_col(df, "SUP_INDEMNIZADA"):
        agg["Ha"] = ("SUP_INDEMNIZADA", "sum")
    if _safe_col(df, "N_PRODUCTORES"):
        df = df.copy()
        df["_PROD_BENEF"] = pd.to_numeric(df["N_PRODUCTORES"], errors="coerce").fillna(0)
        if _safe_col(df, "INDEMNIZACION"):
            _ind = pd.to_numeric(df["INDEMNIZACION"], errors="coerce").fillna(0)
            df["_PROD_BENEF"] = df["_PROD_BENEF"].where(_ind > 0, 0)
        agg["Productores"] = ("_PROD_BENEF", "sum")
    result = df.groupby(col).agg(**agg).reset_index()
    result = result.sort_values("Avisos", ascending=False).head(n)
    rows = []
    for _, r in result.iterrows():
        row = {"name": str(r[col]), "avisos": int(r["Avisos"])}
        if "Indemnización" in r:
            row["indem"] = float(r["Indemnización"])
        if "Desembolso" in r:
            row["desemb"] = float(r["Desembolso"])
        if "Ha" in r:
            row["ha"] = round(float(r["Ha"]), 2)
        if "Productores" in r:
            row["prod"] = int(r["Productores"])
        rows.append(row)
    return rows


def _tipo_breakdown(df):
    """Distribución por tipo de siniestro."""
    if not _safe_col(df, "TIPO_SINIESTRO"):
        return []
    agg = {"Avisos": ("TIPO_SINIESTRO", "count")}
    if _safe_col(df, "INDEMNIZACION"):
        agg["Indemnización"] = ("INDEMNIZACION", "sum")
    result = df.groupby("TIPO_SINIESTRO").agg(**agg).reset_index()
    result = result.sort_values("Avisos", ascending=False)
    rows = []
    for _, r in result.iterrows():
        row = {"tipo": str(r["TIPO_SINIESTRO"]), "avisos": int(r["Avisos"])}
        if "Indemnización" in r:
            row["indem"] = float(r["Indemnización"])
        rows.append(row)
    return rows


def _empresa_breakdown(df):
    """Breakdown por empresa."""
    if not _safe_col(df, "EMPRESA"):
        return []
    results = []
    for emp in df["EMPRESA"].dropna().unique():
        df_emp = df[df["EMPRESA"] == emp]
        m = _calcular_metricas(df_emp)
        results.append({"empresa": str(emp), **m})
    return results


def _dictamen_breakdown(df):
    """Breakdown por resultado de dictamen."""
    for col in ["DICTAMEN", "RESULTADO_AJUSTE", "RESULTADO_INSPECCION"]:
        if _safe_col(df, col):
            counts = df[col].astype(str).str.upper().value_counts().to_dict()
            return {k: int(v) for k, v in counts.items() if k not in ("NAN", "NONE", "")}
    return {}


LLUVIA_TYPES = {"INUNDACION", "INUNDACIÓN", "HUAYCO", "HUAICO",
                "LLUVIAS EXCESIVAS", "DESLIZAMIENTO", "DESLIZAMIENTOS"}


def _generar_insights(df, metricas, tipos, provincias_o_distritos=None,
                      col_geo="PROVINCIA", provs_seleccionadas=None):
    """Genera insights automáticos basados en los datos."""
    insights = []
    m = metricas
    n = m["avisos"]
    if n == 0:
        return insights

    if tipos and len(tipos) > 0:
        top = tipos[0]
        pct = (top["avisos"] / n * 100) if n > 0 else 0
        indem_txt = f" y {_fmt_money_py(top.get('indem', 0))} en indemnización" if top.get("indem") else ""
        insights.append({
            "title": f"{top['tipo']} es el siniestro predominante",
            "text": f"con {top['avisos']:,} avisos ({pct:.1f}% del total){indem_txt}.",
            "type": "predominance"
        })

    if tipos:
        lluvia = [t for t in tipos if t["tipo"].upper() in LLUVIA_TYPES]
        if lluvia:
            lluvia_avisos = sum(t["avisos"] for t in lluvia)
            lluvia_indem = sum(t.get("indem", 0) for t in lluvia)
            lluvia_nombres = ", ".join(t["tipo"].lower() for t in lluvia[:3])
            insights.append({
                "title": "Eventos asociados a lluvias",
                "text": f"({lluvia_nombres}) suman {lluvia_avisos:,} avisos y {_fmt_money_py(lluvia_indem)}.",
                "type": "lluvia"
            })

    if provincias_o_distritos:
        rezago = [p for p in provincias_o_distritos
                  if p["avisos"] >= 5 and p.get("indem", 0) == 0]
        if rezago:
            top_rez = rezago[0]
            insights.append({
                "title": f"{top_rez['name']} presenta rezago en evaluación",
                "text": f"con {top_rez['avisos']} avisos sin indemnización registrada.",
                "type": "rezago"
            })

    if m["pct_desembolso"] > 0:
        if m["pct_desembolso"] < 30:
            insights.append({
                "title": "Bajo nivel de desembolso",
                "text": f"Solo {m['pct_desembolso']:.1f}% de la indemnización ha sido desembolsada ({_fmt_money_py(m['desembolso'])}).",
                "type": "alert"
            })
        elif m["pct_desembolso"] >= 90:
            insights.append({
                "title": "Alto nivel de desembolso",
                "text": f"{m['pct_desembolso']:.1f}% de la indemnización ya fue desembolsada.",
                "type": "positive"
            })

    if provs_seleccionadas and provincias_o_distritos:
        sel = [p for p in provincias_o_distritos if p["name"] in provs_seleccionadas]
        if sel:
            combined_avisos = sum(p["avisos"] for p in sel)
            combined_indem = sum(p.get("indem", 0) for p in sel)
            combined_ha = sum(p.get("ha", 0) for p in sel)
            combined_prod = sum(p.get("prod", 0) for p in sel)
            nombres = " + ".join(p["name"] for p in sel)
            insights.append({
                "title": f"{nombres}: foco seleccionado",
                "text": f"{combined_avisos:,} avisos combinados, {_fmt_money_py(combined_indem)} indemnización, {combined_ha:,.1f} ha, {combined_prod:,} productores.",
                "type": "highlight"
            })

    return insights[:4]


def _empresa_composition(df):
    """Describe composición de empresa."""
    if not _safe_col(df, "EMPRESA"):
        return ""
    counts = df["EMPRESA"].value_counts()
    total = counts.sum()
    if len(counts) == 1:
        return f"Opera exclusivamente con {counts.index[0]}"
    parts = []
    for emp, cnt in counts.items():
        pct = cnt / total * 100
        parts.append(f"{emp}: {pct:.0f}%")
    return " · ".join(parts)


def _fmt_money_py(n):
    """Format money in Python."""
    if n is None or n == 0:
        return "S/ 0"
    if abs(n) >= 1_000_000:
        return f"S/ {n/1_000_000:,.2f} M"
    return f"S/ {n:,.0f}"


# ══════════════════════════════════════════════════════════════════
# PREPARAR DATA
# ══════════════════════════════════════════════════════════════════

def _prepare_data(df, filtros, fecha_corte):
    """Prepara toda la data en un dict para la generación de slides.
    Soporta modelo acumulativo: Nacional + Departamental + Provincial + Distrital.
    """
    scope = filtros.get("scope", "nacional")
    incluir_nacional = filtros.get("incluir_nacional", True)
    deptos = filtros.get("departamentos", [])
    provs = filtros.get("provincias", [])
    dists = filtros.get("distritos", [])

    filtros_base = {k: v for k, v in filtros.items()
                    if k not in ("departamentos", "provincias", "distritos", "scope", "incluir_nacional")}
    df_base = _aplicar_filtros(df, filtros_base)

    data = {
        "fecha_corte": fecha_corte,
        "scope": scope,
        "filtros": {
            "deptos": deptos, "provs": provs, "dists": dists,
            "tipos": filtros.get("tipos_siniestro", []),
            "empresa": filtros.get("empresa", "ambas"),
            "fecha_inicio": str(filtros.get("fecha_inicio", "")) if filtros.get("fecha_inicio") else "",
            "fecha_fin": str(filtros.get("fecha_fin", "")) if filtros.get("fecha_fin") else "",
        },
        "sections": [],
    }

    if incluir_nacional or not deptos:
        m = _calcular_metricas(df_base)
        tipos_nac = _tipo_breakdown(df_base)
        top_deptos_nac = _top_breakdown(df_base, "DEPARTAMENTO", 10)
        n_deptos = df_base["DEPARTAMENTO"].nunique() if _safe_col(df_base, "DEPARTAMENTO") else 0
        data["sections"].append({
            "type": "nacional",
            "metricas": m,
            "pipeline": _calcular_pipeline(df_base),
            "dictamen": _dictamen_breakdown(df_base),
            "empresas": _empresa_breakdown(df_base),
            "top_deptos": top_deptos_nac,
            "tipos": tipos_nac,
            "n_deptos": n_deptos,
            "insights": _generar_insights(df_base, m, tipos_nac, top_deptos_nac, "DEPARTAMENTO"),
        })

    if deptos:
        for depto in deptos:
            df_d = df_base[df_base["DEPARTAMENTO"] == depto] if _safe_col(df_base, "DEPARTAMENTO") else pd.DataFrame()
            if len(df_d) == 0:
                continue
            m = _calcular_metricas(df_d)
            tipos_d = _tipo_breakdown(df_d)
            provs_d = _top_breakdown(df_d, "PROVINCIA", 20)
            emp_comp = _empresa_composition(df_d)
            n_provs = df_d["PROVINCIA"].nunique() if _safe_col(df_d, "PROVINCIA") else 0
            data["sections"].append({
                "type": "departamental",
                "name": depto,
                "metricas": m,
                "pipeline": _calcular_pipeline(df_d),
                "dictamen": _dictamen_breakdown(df_d),
                "provincias": provs_d,
                "tipos": tipos_d,
                "empresa_comp": emp_comp,
                "n_provincias": n_provs,
                "insights": _generar_insights(df_d, m, tipos_d, provs_d, "PROVINCIA", provs),
                "provs_seleccionadas": provs,
            })

    if provs:
        for prov in provs:
            df_p = df_base[df_base["PROVINCIA"] == prov] if _safe_col(df_base, "PROVINCIA") else pd.DataFrame()
            if len(df_p) == 0:
                continue
            m = _calcular_metricas(df_p)
            depto_name = str(df_p["DEPARTAMENTO"].iloc[0]) if _safe_col(df_p, "DEPARTAMENTO") and len(df_p) > 0 else ""
            tipos_p = _tipo_breakdown(df_p)
            dists_p = _top_breakdown(df_p, "DISTRITO", 20)
            emp_comp = _empresa_composition(df_p)
            data["sections"].append({
                "type": "provincial",
                "name": prov,
                "depto": depto_name,
                "metricas": m,
                "pipeline": _calcular_pipeline(df_p),
                "dictamen": _dictamen_breakdown(df_p),
                "distritos": dists_p,
                "tipos": tipos_p,
                "empresa_comp": emp_comp,
                "insights": _generar_insights(df_p, m, tipos_p, dists_p, "DISTRITO"),
            })

    if dists:
        for dist in dists[:5]:
            df_dist = df_base[df_base["DISTRITO"] == dist] if _safe_col(df_base, "DISTRITO") else pd.DataFrame()
            if len(df_dist) == 0:
                continue
            m = _calcular_metricas(df_dist)
            prov_name = str(df_dist["PROVINCIA"].iloc[0]) if _safe_col(df_dist, "PROVINCIA") and len(df_dist) > 0 else ""
            depto_name = str(df_dist["DEPARTAMENTO"].iloc[0]) if _safe_col(df_dist, "DEPARTAMENTO") and len(df_dist) > 0 else ""
            data["sections"].append({
                "type": "distrital",
                "name": dist,
                "prov": prov_name,
                "depto": depto_name,
                "metricas": m,
                "pipeline": _calcular_pipeline(df_dist),
                "tipos": _tipo_breakdown(df_dist),
            })

    return data


# ══════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS FOR PPT GENERATION (python-pptx)
# ══════════════════════════════════════════════════════════════════

# Color palette
C = {
    "forest": RGBColor(0x1B, 0x43, 0x32),
    "green": RGBColor(0x2D, 0x6A, 0x4F),
    "sage": RGBColor(0x52, 0xB7, 0x88),
    "mint": RGBColor(0x95, 0xD5, 0xB2),
    "cream": RGBColor(0xF5, 0xF1, 0xEB),
    "gold": RGBColor(0xD4, 0xA8, 0x43),
    "amber": RGBColor(0xC1, 0x78, 0x17),
    "navy": RGBColor(0x1A, 0x27, 0x44),
    "dark": RGBColor(0x21, 0x25, 0x29),
    "gray": RGBColor(0x6C, 0x75, 0x7D),
    "lightGray": RGBColor(0xE9, 0xEC, 0xEF),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "red": RGBColor(0xC0, 0x39, 0x2B),
    "blue": RGBColor(0x21, 0x96, 0xF3),
    "orange": RGBColor(0xE6, 0x7E, 0x22),
    "coral": RGBColor(0xE7, 0x4C, 0x3C),
    "teal": RGBColor(0x2B, 0xA5, 0xA5),
    "yellowBg": RGBColor(0xFF, 0xF3, 0xCD),
    "lightCream": RGBColor(0xF8, 0xF9, 0xFA),
}


def _fmt_num(n):
    """Format number with locale."""
    if n is None:
        return "0"
    return f"{int(n):,}"


def _fmt_money(n):
    """Format money (S/ X,XXX or S/ X.XX M)."""
    if n is None or n == 0:
        return "S/ 0"
    if abs(n) >= 1_000_000:
        return f"S/{n/1_000_000:,.2f}M"
    if abs(n) >= 100_000:
        return f"S/{n/1_000:,.0f}K"
    return f"S/ {n:,.0f}"


def _fmt_pct(n):
    """Format percentage."""
    if n is None:
        return "0%"
    return f"{n:.1f}%"


def _add_shadow(shape, alpha_val=15000):
    """Add shadow effect to shape via XML."""
    try:
        spPr = shape._element.spPr
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
            'blurRad': '50800',
            'dist': '25400',
            'dir': '8100000',
        })
        srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_val)})
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)
    except Exception:
        pass


def _make_logo(slide, x, y, w, h):
    """Create MIDAGRI/SAC logo using shapes: white rounded rect with two teal rects inside."""
    logo_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, w, h
    )
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = C["white"]
    logo_bg.line.color.rgb = C["teal"]
    logo_bg.line.width = Pt(2)

    rect_w = w * 0.35
    rect_h = h * 0.6
    gap = (w - rect_w * 2) / 3

    rect1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + gap, y + (h - rect_h) / 2,
        rect_w, rect_h
    )
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = C["teal"]
    rect1.line.fill.background()

    rect2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + gap * 2 + rect_w, y + (h - rect_h) / 2,
        rect_w, rect_h
    )
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = C["sage"]
    rect2.line.fill.background()


def _add_kpi_card(slide, left, top, w, h, label, value, sublabel, accent_color, icon_text):
    """Add enhanced metric card with colored top bar, icon circle, and large number."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, w, h
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["white"]
    bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    bg.line.width = Pt(1)
    _add_shadow(bg, 15000)

    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, w, Inches(0.25)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()

    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + (w - Inches(0.45)) / 2, top + Inches(0.35),
        Inches(0.45), Inches(0.45)
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = accent_color
    icon_circle.line.fill.background()

    icon_tf = slide.shapes.add_textbox(
        left + (w - Inches(0.45)) / 2, top + Inches(0.35),
        Inches(0.45), Inches(0.45)
    )
    icon_frame = icon_tf.text_frame
    icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    icon_p = icon_frame.paragraphs[0]
    icon_p.text = icon_text
    icon_p.font.size = Pt(20)
    icon_p.alignment = PP_ALIGN.CENTER

    # Auto-size value font based on text length
    val_str = str(value)
    if len(val_str) > 10:
        val_font = Pt(18)
    elif len(val_str) > 7:
        val_font = Pt(22)
    else:
        val_font = Pt(26)

    value_tf = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(0.88),
        w - Inches(0.16), Inches(0.45)
    )
    value_frame = value_tf.text_frame
    value_frame.word_wrap = False
    value_frame.margin_top = Pt(0)
    value_frame.margin_bottom = Pt(0)
    value_p = value_frame.paragraphs[0]
    value_p.text = val_str
    value_p.font.size = val_font
    value_p.font.bold = True
    value_p.font.color.rgb = C["navy"]
    value_p.font.name = "Georgia"
    value_p.alignment = PP_ALIGN.CENTER

    sub_tf = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(1.35),
        w - Inches(0.16), Inches(0.25)
    )
    sub_frame = sub_tf.text_frame
    sub_frame.word_wrap = True
    sub_frame.margin_top = Pt(0)
    sub_frame.margin_bottom = Pt(0)
    sub_p = sub_frame.paragraphs[0]
    sub_p.text = str(sublabel) if sublabel else ""
    sub_p.font.size = Pt(8)
    sub_p.font.color.rgb = C["gray"]
    sub_p.alignment = PP_ALIGN.CENTER

    label_tf = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(1.58),
        w - Inches(0.16), Inches(0.22)
    )
    label_frame = label_tf.text_frame
    label_frame.word_wrap = True
    label_frame.margin_top = Pt(0)
    label_frame.margin_bottom = Pt(0)
    label_p = label_frame.paragraphs[0]
    label_p.text = str(label)
    label_p.font.size = Pt(9)
    label_p.font.color.rgb = C["gray"]
    label_p.alignment = PP_ALIGN.CENTER


def _add_header_bar(slide, title, color, y_pos=Inches(0.3)):
    """Add colored rectangle header bar with title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), y_pos,
        Inches(10), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.alignment = PP_ALIGN.LEFT
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(1.2),
                     col_widths=None, has_total=False, max_rows=12):
    """Add professional table with alternating rows, header styling, optional total row."""
    rows_to_add = min(len(rows), max_rows)
    if has_total:
        rows_to_add += 1
    cols = len(headers)

    table_shape = slide.shapes.add_table(rows_to_add + 1, cols, left, top,
                                         Inches(9.4), Inches(0.35 * (rows_to_add + 1)))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["navy"]
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C["white"]
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows[:max_rows], 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = C["white"]
            else:
                cell.fill.fore_color.rgb = C["lightCream"]

            p = cell.text_frame.paragraphs[0]
            p.text = str(value) if value is not None else ""
            p.font.size = Pt(9)
            p.font.color.rgb = C["dark"]
            p.alignment = PP_ALIGN.CENTER

    if has_total and len(rows) > 0:
        total_row = rows_to_add
        for col_idx in range(cols):
            cell = table.cell(total_row, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["navy"]
            p = cell.text_frame.paragraphs[0]
            p.text = "TOTAL" if col_idx == 0 else ""
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = C["white"]
            p.alignment = PP_ALIGN.CENTER


def _add_pipeline_flow(slide, pipeline, y_pos):
    """Add 5-step pipeline flow with numbered circles and percentages."""
    if not pipeline:
        return

    total = sum(p["val"] for p in pipeline)
    x_start = Inches(0.4)
    card_w = Inches(1.7)
    card_h = Inches(1.8)
    circle_size = Inches(0.7)

    colors = [C["teal"], C["amber"], C["teal"], C["navy"], C["orange"]]

    for i, stage in enumerate(pipeline[:5]):
        x_pos = x_start + Inches(i * 1.85)

        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, card_w, card_h
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = C["white"]
        card_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        card_bg.line.width = Pt(1)
        _add_shadow(card_bg, 12000)

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x_pos, y_pos, card_w, Inches(0.18)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors[i]
        top_bar.line.fill.background()

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x_pos + (card_w - circle_size) / 2, y_pos + Inches(0.25),
            circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()

        circle_tf = slide.shapes.add_textbox(
            x_pos + (card_w - circle_size) / 2, y_pos + Inches(0.25),
            circle_size, circle_size
        )
        circle_frame = circle_tf.text_frame
        circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        circle_p = circle_frame.paragraphs[0]
        circle_p.text = str(i + 1)
        circle_p.font.size = Pt(24)
        circle_p.font.bold = True
        circle_p.font.color.rgb = C["white"]
        circle_p.alignment = PP_ALIGN.CENTER

        label_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.05), y_pos + Inches(1.05),
            card_w - Inches(0.1), Inches(0.3)
        )
        label_frame = label_tf.text_frame
        label_frame.word_wrap = True
        label_p = label_frame.paragraphs[0]
        label_p.text = stage["label"]
        label_p.font.size = Pt(9)
        label_p.font.color.rgb = C["gray"]
        label_p.alignment = PP_ALIGN.CENTER

        val_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.05), y_pos + Inches(1.35),
            card_w - Inches(0.1), Inches(0.28)
        )
        val_frame = val_tf.text_frame
        val_frame.word_wrap = True
        val_p = val_frame.paragraphs[0]
        val_p.text = f"{stage['val']:,}"
        val_p.font.size = Pt(16)
        val_p.font.bold = True
        val_p.font.color.rgb = C["navy"]
        val_p.alignment = PP_ALIGN.CENTER

        pct = (stage["val"] / total * 100) if total > 0 else 0
        pct_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.05), y_pos + Inches(1.63),
            card_w - Inches(0.1), Inches(0.15)
        )
        pct_frame = pct_tf.text_frame
        pct_frame.word_wrap = True
        pct_p = pct_frame.paragraphs[0]
        pct_p.text = f"{pct:.1f}%"
        pct_p.font.size = Pt(10)
        pct_p.font.bold = True
        pct_p.font.color.rgb = colors[i]
        pct_p.alignment = PP_ALIGN.CENTER

        if i < len(pipeline) - 1 and i < 4:
            dash_x = x_pos + card_w + Inches(0.05)
            dash_y = y_pos + Inches(0.9)
            for j in range(3):
                dash = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    dash_x + Inches(j * 0.08), dash_y,
                    Inches(0.06), Inches(0.1)
                )
                dash.fill.solid()
                dash.fill.fore_color.rgb = C["lightGray"]
                dash.line.fill.background()


def _add_alert_box(slide, text, left, top, width):
    """Add yellow/amber alert box with warning content."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, Inches(0.55)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["yellowBg"]
    bg.line.color.rgb = C["amber"]
    bg.line.width = Pt(1.5)

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, Inches(0.05), Inches(0.55)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = C["amber"]
    accent_bar.line.fill.background()

    tf = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.08),
        width - Inches(0.3), Inches(0.4)
    )
    text_frame = tf.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "⚠ " + text
    p.font.size = Pt(10)
    p.font.color.rgb = C["dark"]
    p.alignment = PP_ALIGN.LEFT


def _add_resumen_ejecutivo_slide(slide, prs, dept_name, text, fecha_corte):
    """Add dark background resumen ejecutivo slide with narrative text."""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.2),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.42),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.4),
        Inches(2.0), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["teal"]
    badge.line.width = Pt(1)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.4),
        Inches(2.0), Inches(0.35)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = dept_name.upper()
    badge_p.font.size = Pt(12)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["teal"]
    badge_p.alignment = PP_ALIGN.CENTER

    title_tf = slide.shapes.add_textbox(
        Inches(2.6), Inches(0.35),
        Inches(6.8), Inches(0.45)
    )
    title_frame = title_tf.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "Resumen Ejecutivo"
    title_p.font.size = Pt(24)
    title_p.font.bold = True
    title_p.font.color.rgb = C["lightGray"]
    title_p.font.name = "Georgia"
    title_p.alignment = PP_ALIGN.LEFT

    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(0.85),
        Inches(9.2), Inches(0.02)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = C["teal"]
    sep_line.line.fill.background()

    body_tf = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.1),
        Inches(8.8), Inches(4.15)
    )
    body_frame = body_tf.text_frame
    body_frame.word_wrap = True
    body_p = body_frame.paragraphs[0]
    body_p.text = text
    body_p.font.size = Pt(15)
    body_p.font.color.rgb = C["white"]
    body_p.font.name = "Georgia"
    body_p.alignment = PP_ALIGN.JUSTIFY
    body_p.line_spacing = 1.4

    footer_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.2),
        Inches(9.2), Inches(0.2)
    )
    footer_frame = footer_tf.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = f"Fuente: DSFFA — MIDAGRI, SAC 2025-2026. Corte al {fecha_corte}."
    footer_p.font.size = Pt(9)
    footer_p.font.italic = True
    footer_p.font.color.rgb = C["lightGray"]
    footer_p.alignment = PP_ALIGN.RIGHT


# ══════════════════════════════════════════════════════════════════
# SLIDE GENERATION FUNCTIONS
# ══════════════════════════════════════════════════════════════════

def _generar_resumen_texto(section_data, scope_type):
    """Generate narrative paragraph text from section metrics."""
    m = section_data.get("metricas", {})
    tipos = section_data.get("tipos", [])
    provs = section_data.get("provincias", [])

    nombre = section_data.get("name", "La región")
    avisos = m.get("avisos", 0)
    pct_eval = m.get("pct_eval", 0)
    cerrados = m.get("cerrados", 0)
    indem_total = m.get("indemnizacion", 0)
    pct_desembolso = m.get("pct_desembolso", 0)
    desembolso = m.get("desembolso", 0)

    tipo_predominante = tipos[0]["tipo"] if tipos else "datos"
    prov_principal = provs[0]["name"] if provs else "principales distritos"

    text = (
        f"{nombre} registra {avisos:,} avisos con avance de evaluación de {pct_eval:.1f}% "
        f"({cerrados:,} expedientes cerrados). La indemnización total reconocida asciende a "
        f"S/ {indem_total:,.0f}, con un desembolso acumulado de S/ {desembolso:,.0f} "
        f"({pct_desembolso:.1f}% del monto reconocido). El tipo de siniestro predominante es "
        f"{tipo_predominante}, concentrándose la mayor carga en {prov_principal}. "
        f"Se observa rezago en la evaluación de expedientes pendientes, requiriendo acelerar "
        f"los procesos de inspección y dictamen para cumplir los compromisos de desembolso. "
        f"Las principales causas de atraso incluyen limitaciones en recursos humanos, "
        f"complejidad de los siniestros y coordinación interinstitucional."
    )
    return text[:200]  # Limit to ~200 words


def _add_pipeline_slide(prs, pipeline, dictamen, metricas):
    """Add 5-step pipeline process flow slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    tf_title = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(9.2), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Flujo de Procesos del SAC — Estado Actual"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    _add_pipeline_flow(slide, pipeline, Inches(1.1))

    info_text = "⚠ Avisos dictaminados NO INDEMNIZABLE: requieren documentación adicional"
    _add_alert_box(slide, info_text, Inches(0.4), Inches(4.9), Inches(9.2))


def _add_resumen_ejecutivo(prs, section_name, resumen_text, scope_label, fecha_corte):
    """Add resumen ejecutivo slide with dark background and narrative text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_resumen_ejecutivo_slide(slide, prs, section_name, resumen_text, fecha_corte)


def _add_tipo_siniestro_slide(prs, tipos):
    """Add tipo de siniestro distribution slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    tf_title = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(9.2), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Distribución por Tipo de Siniestro"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    headers = ["Tipo Siniestro", "Monto (S/)", "%"]
    rows = []
    total_monto = sum(t.get("indem", 0) for t in tipos)

    for t in tipos[:12]:
        monto = t.get("indem", 0)
        pct = (monto / total_monto * 100) if total_monto > 0 else 0
        rows.append([
            t["tipo"],
            _fmt_money(monto),
            f"{pct:.1f}%"
        ])

    col_widths = [Inches(3.5), Inches(3.0), Inches(1.2)]
    _add_styled_table(slide, headers, rows, left=Inches(1.0), top=Inches(1.3),
                     col_widths=col_widths, max_rows=12)

    footer_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.35),
        Inches(9.2), Inches(0.2)
    )
    footer_frame = footer_tf.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Fuente: DSFFA — MIDAGRI, SAC 2025-2026"
    footer_p.font.size = Pt(8)
    footer_p.font.italic = True
    footer_p.font.color.rgb = C["gray"]
    footer_p.alignment = PP_ALIGN.RIGHT


def _add_top_deptos_chart(prs, top_deptos):
    """Add top departamentos bar chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    tf_title = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(9.2), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Top 12 Departamentos por Indemnización Reconocida"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = [d["name"] for d in top_deptos[:12]]

    values = [d.get("indem", 0) for d in top_deptos[:12]]
    chart_data.add_series('Indemnización (S/)', tuple(values))

    # Add chart
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.3),
        Inches(9.0), Inches(3.8), chart_data
    ).chart

    chart_shape.has_legend = False
    chart_shape.series[0].format.fill.solid()
    chart_shape.series[0].format.fill.fore_color.rgb = C["navy"]

    footer_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.35),
        Inches(9.2), Inches(0.2)
    )
    footer_frame = footer_tf.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Fuente: DSFFA — MIDAGRI, SAC 2025-2026"
    footer_p.font.size = Pt(8)
    footer_p.font.italic = True
    footer_p.font.color.rgb = C["gray"]
    footer_p.alignment = PP_ALIGN.RIGHT


def _add_portada(prs, data):
    """Add cover slide with dark navy background and teal accents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    _make_logo(slide, Inches(4.0), Inches(0.55), Inches(2.0), Inches(1.6))

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.3),
        Inches(9), Inches(0.55)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    p = text_frame.paragraphs[0]
    p.text = "SEGURO AGRÍCOLA CATASTRÓFICO"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    tf_subtitle = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.95),
        Inches(9), Inches(0.35)
    )
    text_frame = tf_subtitle.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "SAC 2025–2026"
    p.font.size = Pt(20)
    p.font.color.rgb = C["teal"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(3.45),
        Inches(3), Inches(0.03)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = C["teal"]
    sep_line.line.fill.background()

    scope_text = "Resumen de Avance Nacional"
    if data.get("filtros", {}).get("deptos"):
        scope_text = f"Departamento: {', '.join(data['filtros']['deptos'][:2])}"

    tf_scope = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.6),
        Inches(9), Inches(0.3)
    )
    text_frame = tf_scope.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = scope_text
    p.font.size = Pt(14)
    p.font.color.rgb = C["white"]
    p.alignment = PP_ALIGN.CENTER

    fecha_corte = data.get("fecha_corte", "S.F.")
    tf_footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.85),
        Inches(9), Inches(0.35)
    )
    text_frame = tf_footer.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Dirección de Seguro y Fomento del Financiamiento Agrario — MIDAGRI"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = C["lightGray"]
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = f"Corte al {fecha_corte}"
    p2.font.size = Pt(9)
    p2.font.italic = True
    p2.font.color.rgb = C["lightGray"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)


def _add_cierre(prs, fecha_corte):
    """Add closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    _make_logo(slide, Inches(4.0), Inches(0.55), Inches(2.0), Inches(1.3))

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(9), Inches(0.55)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    p = text_frame.paragraphs[0]
    p.text = "SEGURO AGRÍCOLA CATASTRÓFICO"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    tf_subtitle = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.65),
        Inches(9), Inches(0.35)
    )
    text_frame = tf_subtitle.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "SAC 2025–2026"
    p.font.size = Pt(20)
    p.font.color.rgb = C["teal"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(3.05),
        Inches(3), Inches(0.03)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = C["teal"]
    sep_line.line.fill.background()

    tf_footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.4),
        Inches(9), Inches(1.0)
    )
    text_frame = tf_footer.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Dirección de Seguro y Fomento del Financiamiento Agrario"
    p.font.size = Pt(13)
    p.font.color.rgb = C["white"]
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = "Ministerio de Desarrollo Agrario y Riego — MIDAGRI"
    p2.font.size = Pt(12)
    p2.font.color.rgb = C["lightGray"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    p3 = text_frame.add_paragraph()
    p3.text = f"Corte al {fecha_corte}"
    p3.font.size = Pt(10)
    p3.font.italic = True
    p3.font.color.rgb = C["lightGray"]
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)


def _add_nacional_section(prs, section):
    """Add nacional section: indicadores clave, pipeline, tipo siniestro, top deptos."""
    m = section["metricas"]
    empresas = section.get("empresas", [])
    pipeline = section.get("pipeline", [])
    tipos = section.get("tipos", [])
    deptos = section.get("departamentos", [])

    # Slide 1: Indicadores Clave
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    _make_logo(slide, Inches(0.4), Inches(0.35), Inches(0.6), Inches(0.65))

    tf_title = slide.shapes.add_textbox(
        Inches(1.2), Inches(0.3),
        Inches(8.3), Inches(0.55)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Resumen Nacional — Indicadores Clave"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.LEFT

    tf_sub = slide.shapes.add_textbox(
        Inches(1.2), Inches(0.85),
        Inches(8.3), Inches(0.25)
    )
    text_frame = tf_sub.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    emp_text = "Consolidado La Positiva + Rímac · Campaña 2025-2026"
    if empresas:
        emp_text = " + ".join([f"{e['empresa']} ({e['avisos']} avisos)" for e in empresas[:2]])
    p.text = emp_text
    p.font.size = Pt(10)
    p.font.color.rgb = C["gray"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.15),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    kpi_configs = [
        ("Avisos Reportados", _fmt_num(m["avisos"]), "24 departamentos", C["teal"], "⚠"),
        ("Avance Evaluación", f"{_fmt_pct(m['pct_eval'])}", f"{m['cerrados']:,} cerrados de {m['avisos']:,}", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "Reconocida a productores", C["orange"], "●"),
        ("Avance Desembolso", f"{_fmt_pct(m['pct_desembolso'])}", f"{_fmt_money(m['desembolso'])} desembolsados", C["teal"], "✦"),
        ("Ha Indemnizadas", f"{m['ha_indemnizadas']:,.1f}", "Solo con evaluación cerrada", C["teal"], "🌿"),
        ("Productores", _fmt_num(m["productores"]), "Beneficiados con indemnización", C["sage"], "👥"),
    ]

    if empresas:
        for i, emp in enumerate(empresas[:2]):
            emp_label = f"{emp['empresa']}"
            emp_val = f"{_fmt_num(emp['avisos'])}"
            emp_detail = f"{_fmt_pct(emp['pct_eval'])} cerr. · {_fmt_money(emp['indemnizacion'])}"
            kpi_configs.append((emp_label, emp_val, emp_detail, C["teal"], f"{i+1}"))

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs[:8]):
        col = i % 4
        row = i // 4
        left = Inches(0.35 + col * 2.35)
        top = Inches(1.65 + row * 2.0)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    # Slide 2: Pipeline de Procesos
    if pipeline:
        _add_pipeline_slide(prs, pipeline, {}, m)

    # Slide 3: Tipo de Siniestro
    if tipos:
        _add_tipo_siniestro_slide(prs, tipos)

    # Slide 4: Top Departamentos
    if deptos:
        _add_top_deptos_chart(prs, deptos)


def _add_departamental_section(prs, section, fecha_corte="S.F."):
    """Add departamental section: separator + metrics + resumen ejecutivo."""
    name = section.get("name", "Departamento")
    m = section["metricas"]
    tipos = section.get("tipos", [])
    provs = section.get("provincias", [])
    pipeline = section.get("pipeline", [])
    insights = section.get("insights", [])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["teal"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "🌾"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7),
        Inches(9), Inches(0.8)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = name.upper()
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = f"{m['avisos']:,} avisos — {_fmt_money(m['indemnizacion'])} indemnización"
    p2.font.size = Pt(13)
    p2.font.color.rgb = C["teal"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.35),
        Inches(1.5), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["orange"]
    badge.line.width = Pt(1.5)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(1.5), Inches(0.35)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = name.upper()
    badge_p.font.size = Pt(11)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["orange"]
    badge_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(2.1), Inches(0.28),
        Inches(7.5), Inches(0.5)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Análisis Departamental — Indicadores Clave"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    kpi_configs = [
        ("Avisos Reportados", _fmt_num(m["avisos"]), "", C["amber"], "⚠"),
        ("Avance Evaluación", f"{_fmt_pct(m['pct_eval'])}", f"{m['cerrados']:,} cerrados", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "", C["orange"], "●"),
        ("Avance Desembolso", f"{_fmt_pct(m['pct_desembolso'])}", _fmt_money(m["desembolso"]), C["teal"], "✦"),
    ]

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs):
        left = Inches(0.3 + i * 2.35)
        top = Inches(1.25)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    if tipos:
        headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)"]
        rows = []
        for t in tipos[:8]:
            rows.append([
                t["tipo"],
                _fmt_num(t["avisos"]),
                _fmt_money(t.get("indem", 0)),
            ])
        col_widths = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(3.25), col_widths=col_widths, max_rows=5)

    if provs:
        headers_p = ["Provincia", "Avisos", "Indem. (S/)"]
        rows_p = []
        for p in provs[:5]:
            rows_p.append([
                p["name"],
                _fmt_num(p["avisos"]),
                _fmt_money(p.get("indem", 0)),
            ])
        col_widths_p = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers_p, rows_p, left=Inches(5.0), top=Inches(3.25), col_widths=col_widths_p, max_rows=5)

    # Slide 3: Resumen Ejecutivo
    resumen_text = _generar_resumen_texto(section, "departamental")
    _add_resumen_ejecutivo(prs, name, resumen_text, "Departamental", fecha_corte)


def _add_provincial_section(prs, section):
    """Add provincial section: separator + metrics."""
    name = section.get("name", "Provincia")
    depto = section.get("depto", "")
    m = section["metricas"]
    dists = section.get("distritos", [])
    tipos = section.get("tipos", [])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["forest"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["sage"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["sage"]
    line_bottom.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["sage"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "🏡"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7),
        Inches(9), Inches(0.8)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = name.upper()
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = depto if depto else "Provincia"
    p2.font.size = Pt(13)
    p2.font.color.rgb = C["sage"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)


def _add_distrital_section(prs, section):
    """Add distrital section: separator slide."""
    name = section.get("name", "Distrito")
    prov = section.get("prov", "")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["teal"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "📍"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7),
        Inches(9), Inches(0.8)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = name.upper()
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = f"{prov} Provincia" if prov else "Distrito"
    p2.font.size = Pt(13)
    p2.font.color.rgb = C["teal"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)


# ══════════════════════════════════════════════════════════════════
# MAIN GENERATION FUNCTION
# ══════════════════════════════════════════════════════════════════

def generar_ppt_dinamico(df, filtros, fecha_corte):
    """
    Genera una presentación PPT dinámica con python-pptx.
    Interfaz compatible con app.py: misma firma, misma lógica de datos.

    Args:
        df: DataFrame consolidado (datos["midagri"])
        filtros: dict con selecciones del usuario
        fecha_corte: string con fecha de corte (ej: "14/03/2026")

    Returns:
        bytes del archivo .pptx
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    prs.author = "DSFFA — MIDAGRI"
    prs.title = "SAC 2025-2026 — Presentación Dinámica"

    data = _prepare_data(df, filtros, fecha_corte)

    _add_portada(prs, data)

    for section in data.get("sections", []):
        section_type = section.get("type", "")

        if section_type == "nacional":
            _add_nacional_section(prs, section)
        elif section_type == "departamental":
            _add_departamental_section(prs, section, fecha_corte)
        elif section_type == "provincial":
            _add_provincial_section(prs, section)
        elif section_type == "distrital":
            _add_distrital_section(prs, section)

    _add_cierre(prs, fecha_corte)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()
